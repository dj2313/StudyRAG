import fitz
import docx
import pptx
import easyocr
import base64
import os
import subprocess
from groq import Groq
from ..config import settings

client = Groq(api_key=settings.GROQ_API_KEY)
reader = easyocr.Reader(['en'], gpu=False) # Fallback reader

def is_digital_pdf(path: str) -> bool:
    doc = fitz.open(path)
    text = "".join([page.get_text() for page in doc])
    return len(text.strip()) > 100

def extract_digital_pdf(path: str) -> str:
    out_dir = os.path.dirname(path)
    try:
        subprocess.run(["marker_single", path, "--output_dir", out_dir], check=True, capture_output=True)
        base = os.path.basename(path).replace(".pdf", "")
        md_file = os.path.join(out_dir, base, f"{base}.md")
        if os.path.exists(md_file):
            with open(md_file, "r", encoding="utf-8") as f:
                return f.read()
    except Exception:
        pass
    
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def encode_image(image_path: str) -> str:
    with open(image_path, "rb") as img:
        return base64.b64encode(img.read()).decode('utf-8')

def extract_scanned(path: str) -> str:
    img_path = path
    is_pdf = path.lower().endswith('.pdf')
    if is_pdf:
        doc = fitz.open(path)
        page = doc.load_page(0)
        pix = page.get_pixmap()
        img_path = path + ".png"
        pix.save(img_path)
        
    try:
        base64_image = encode_image(img_path)
        res = client.chat.completions.create(
            model="llama-3.2-90b-vision-preview", # Fallback for llama-4-scout
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all text exactly as written. No extra output."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }],
            max_tokens=1000,
        )
        text = res.choices[0].message.content
    except Exception:
        # Fallback to easyocr
        text = " ".join(reader.readtext(img_path, detail=0))
        
    if is_pdf and os.path.exists(img_path):
        os.remove(img_path)
    return text

def extract_image(path: str) -> str:
    return extract_scanned(path)

def extract_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_pptx(path: str) -> str:
    prs = pptx.Presentation(path)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def route_file(path: str, mime_type: str) -> str:
    ext = path.lower().split('.')[-1]
    if ext == 'pdf' or 'pdf' in mime_type:
        return extract_digital_pdf(path) if is_digital_pdf(path) else extract_scanned(path)
    elif ext in ['png', 'jpg', 'jpeg'] or 'image' in mime_type:
        return extract_image(path)
    elif ext == 'docx':
        return extract_docx(path)
    elif ext == 'pptx':
        return extract_pptx(path)
    elif ext in ['txt', 'md']:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    raise ValueError(f"Unsupported file type: {ext}")
